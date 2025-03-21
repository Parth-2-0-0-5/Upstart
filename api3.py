from fastapi import FastAPI, File, UploadFile
from fastapi.middleware.cors import CORSMiddleware
import fitz  
from pptx import Presentation
import io

app = FastAPI()

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],  
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

def extract_text_from_pdf(file):
    text = ""
    pdf_document = fitz.open(stream=file.read(), filetype="pdf")
    for page in pdf_document:
        text += page.get_text("text") + "\n"
    return text.strip()

def extract_text_from_pptx(file):
    text = ""
    prs = Presentation(io.BytesIO(file.read()))
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text.strip()

@app.post("/upload_pitch/")
async def upload_pitch(file: UploadFile = File(...)):
    if file.filename.endswith(".pdf"):
        extracted_text = extract_text_from_pdf(file)
    elif file.filename.endswith(".pptx"):
        extracted_text = extract_text_from_pptx(file)
    else:
        return {"error": "Unsupported file format. Use PDF or PPTX."}

    analysis_result = analyze_pitch(extracted_text)
    
    

def analyze_pitch(pitch_text):
    return {
        "Overall Score": 85,
        "Section Feedback": {
            "Problem Statement": "Strong",
            "Solution Clarity": "Strong",
            "Market Size": "Needs Detail",
            "Business Model": "Strong",
            "Competition": "Weak",
            "Team": "Strong"
        }
    }
